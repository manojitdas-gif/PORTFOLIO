import os
import sys
import time

def check_dependencies():
    print("Checking python packages...")
    try:
        from PIL import Image, ImageDraw, ImageFont
        print("[OK] Pillow is installed.")
    except ImportError:
        print("[ERROR] Pillow is NOT installed. Installing it now...")
        os.system("python -m pip install pillow")
        from PIL import Image, ImageDraw, ImageFont

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        print("[OK] python-pptx is installed.")
    except ImportError:
        print("[ERROR] python-pptx is NOT installed. Installing it now...")
        os.system("python -m pip install python-pptx")
        from pptx import Presentation

check_dependencies()

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Define paths
DESKTOP_DIR = "C:\\Users\\MONOJIT\\Desktop"
WORKSPACE_DIR = "C:\\Users\\MONOJIT\\Desktop\\stich bio"
FULL_IMG_PATH = os.path.join(DESKTOP_DIR, "stich_bio_full.png")
OUTPUT_PPTX_PATH = os.path.join(WORKSPACE_DIR, "portfolio_presentation.pptx")
OUTPUT_GIF_PATH = os.path.join(WORKSPACE_DIR, "portfolio_walkthrough.gif")

# Check if full page image exists
if not os.path.exists(FULL_IMG_PATH):
    # Fallback to stich_bio_screenshot.png if full is missing
    FULL_IMG_PATH = os.path.join(DESKTOP_DIR, "stich_bio_screenshot.png")
    if not os.path.exists(FULL_IMG_PATH):
        print(f"[ERROR] Screenshot image not found at {DESKTOP_DIR}. Please run screenshot first.")
        sys.exit(1)

print(f"Using image source: {FULL_IMG_PATH}")
img = Image.open(FULL_IMG_PATH)
width, height = img.size
print(f"Image dimensions: {width} x {height}")

# Create slides subdirectory for cropped sections
slides_dir = os.path.join(WORKSPACE_DIR, "presentation_slides")
os.makedirs(slides_dir, exist_ok=True)

# Define crop coordinates (bounding boxes [left, top, right, bottom])
# Since the full page is ~6000px high, let's map out the crop boxes based on sections.
# If height is less than 6000 (e.g. standard screenshot), adjust accordingly.
crops = {}
if height > 4000:
    crops = {
        "01_hero": (0, 0, width, 880),
        "02_about": (0, 880, width, 1600),
        "03_skills": (0, 1600, width, 2400),
        "04_ai_advisor": (0, 2400, width, 3200),
        "05_projects": (0, 3200, width, 4100),
        "06_agency_hub": (0, 4100, width, 4900),
        "07_resume": (0, 4900, width, 5800),
        "08_footer": (0, 5800, width, min(height, 6400))
    }
else:
    # Single viewport fallback
    crops = {
        "01_hero": (0, 0, width, min(height, 800))
    }

cropped_paths = {}
for name, bbox in crops.items():
    t_left, t_top, t_right, t_bottom = bbox
    # Ensure coordinates are within image bounds
    t_bottom = min(t_bottom, height)
    t_top = min(t_top, height)
    if t_bottom - t_top > 100: # only crop if there's height
        cropped = img.crop((t_left, t_top, t_right, t_bottom))
        slide_path = os.path.join(slides_dir, f"{name}.png")
        cropped.save(slide_path, "PNG")
        cropped_paths[name] = slide_path
        print(f"Saved cropped section '{name}' to {slide_path}")

# Sub-apps screenshots
sub_apps = {
    "growthlytics": os.path.join(DESKTOP_DIR, "screenshot_growthlytics.png"),
    "momotantra": os.path.join(DESKTOP_DIR, "screenshot_momotantra.png"),
    "revercafe": os.path.join(DESKTOP_DIR, "screenshot_revercafe.png"),
    "retailflow": os.path.join(DESKTOP_DIR, "screenshot_retailflow.png"),
}

# ----------------- GENERATING POWERPOINT -----------------
print("Generating PowerPoint Presentation...")
prs = Presentation()
# Set slide dimensions to widescreen (16:9)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Scheme
GOLD = RGBColor(201, 168, 76)
DARK_BG = RGBColor(10, 10, 15)
WHITE = RGBColor(240, 238, 252)
GRAY = RGBColor(160, 158, 192)
TEAL = RGBColor(54, 217, 196)
VIOLET = RGBColor(124, 106, 240)

def apply_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

def add_slide_header(slide, title_text, category="PORTFOLIO BRIEFING"):
    # Header container
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = category.upper()
    p0.font.name = 'Courier New'
    p0.font.size = Pt(11)
    p0.font.bold = True
    p0.font.color.rgb = GOLD
    
    p1 = tf.add_paragraph()
    p1.text = title_text
    p1.font.name = 'Arial'
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = WHITE

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[6] # blank layout
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)

# Title Text Box
tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(7.0), Inches(4.5))
tf = tb.text_frame
tf.word_wrap = True

p_sub = tf.paragraphs[0]
p_sub.text = "FOUNDER | OPERATIONS ASSOCIATE | DATA ANALYST"
p_sub.font.name = 'Courier New'
p_sub.font.size = Pt(13)
p_sub.font.bold = True
p_sub.font.color.rgb = GOLD
p_sub.space_after = Pt(20)

p_title = tf.add_paragraph()
p_title.text = "MANOJIT DAS"
p_title.font.name = 'Arial'
p_title.font.size = Pt(56)
p_title.font.bold = True
p_title.font.color.rgb = WHITE
p_title.space_after = Pt(10)

p_desc = tf.add_paragraph()
p_desc.text = "Digital Operations and Analytical Dashboard Portfolio Review.\nBuilt with high-precision metrics, interactive AI integration models, and digital transformation layouts for offline traditional vendors."
p_desc.font.name = 'Arial'
p_desc.font.size = Pt(15)
p_desc.font.color.rgb = GRAY
p_desc.space_after = Pt(30)

p_meta = tf.add_paragraph()
p_meta.text = "Live Portfolio: portfolio-delta-nine-buockq1joj.vercel.app\nLocation: Kolkata, India | Contact: manojitdas876@gmail.com"
p_meta.font.name = 'Courier New'
p_meta.font.size = Pt(11)
p_meta.font.color.rgb = TEAL

# Add hero screenshot on the right
hero_img = cropped_paths.get("01_hero")
if hero_img:
    slide.shapes.add_picture(hero_img, Inches(7.8), Inches(1.2), width=Inches(5.0))

# Slide 2: About & Mindset
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_slide_header(slide, "Operations Mindset & Analytical Precision", "01 -- PROFILE DETAILS")

tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.8), Inches(5.2))
tf = tb.text_frame
tf.word_wrap = True

bullets = [
    ("Core Professional Focus", "Results-oriented Office & Operations Associate and Data Analyst with hands-on experience managing high-volume administrative workflows, data operations, and customer success targets."),
    ("Founder of Growthlytics", "Assists local traditional vendors and offline businesses in West Bengal to go digital, launch online ordering, and automate invoicing."),
    ("Background & Experience", "Managed 50+ daily inbound/outbound customer interactions at Tech Mahindra; completed Business Analytics internships at Bharat Intern; operations coordinator at E-Cell IIT Bombay."),
    ("Core Philosophy", "Bridging the gap between technical data structures (SQL, Power BI, Advanced Excel) and direct customer success and client-facing communication.")
]

for title, body in bullets:
    p = tf.add_paragraph()
    p.text = f"• {title}: "
    p.font.name = 'Arial'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = GOLD
    
    # Append the body text as normal text in the same paragraph
    run = p.add_run()
    run.text = body
    run.font.name = 'Arial'
    run.font.size = Pt(13)
    run.font.bold = False
    run.font.color.rgb = GRAY
    p.space_after = Pt(14)

about_img = cropped_paths.get("02_about")
if about_img:
    slide.shapes.add_picture(about_img, Inches(7.5), Inches(1.6), width=Inches(5.3))

# Slide 3: Competencies
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_slide_header(slide, "Core Professional Competencies", "02 -- EXPERTISE MATRIX")

tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.5), Inches(5.2))
tf = tb.text_frame
tf.word_wrap = True

skills = [
    ("Digital Operations & Automation (85%)", "Process digitization, vendor onboarding, campaign tracking, and operations dashboards."),
    ("Customer Success & SLA Handling (90%)", "CRM tools, SLA tracking, ticket resolution, and customer satisfaction (CSAT) optimization."),
    ("Business Analysis & SQL Databases (80%)", "Writing complex relational queries, data cleansing, structuring data, and Power BI report building."),
    ("AI Integration & Prompt Engineering (70%)", "Leveraging Google Gemini APIs, SDKs, and prompt models to automate text-based business advisories."),
    ("Web Development & Static Layouts (75%)", "React, TypeScript, CSS custom templates, and responsive frontend design."),
    ("Financial Accounting (72%)", "Capital markets, ledger posting, and invoicing (Tally ERP).")
]

for title, body in skills:
    p = tf.add_paragraph()
    p.text = f"• {title}\n"
    p.font.name = 'Arial'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEAL
    
    run = p.add_run()
    run.text = body
    run.font.name = 'Arial'
    run.font.size = Pt(12)
    run.font.bold = False
    run.font.color.rgb = GRAY
    p.space_after = Pt(10)

skills_img = cropped_paths.get("03_skills")
if skills_img:
    slide.shapes.add_picture(skills_img, Inches(7.3), Inches(1.5), width=Inches(5.5))

# Slide 4: Growthlytics AI Advisor
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_slide_header(slide, "Growthlytics AI Problem Solver", "03 -- AI SOLUTIONS")

tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.5), Inches(5.2))
tf = tb.text_frame
tf.word_wrap = True

p_intro = tf.paragraphs[0]
p_intro.text = "Interactive Terminal Simulation:"
p_intro.font.name = 'Courier New'
p_intro.font.size = Pt(16)
p_intro.font.bold = True
p_intro.font.color.rgb = GOLD
p_intro.space_after = Pt(14)

points = [
    ("Goal", "Address the operational bottlenecks of local traditional shops converting to digital formats."),
    ("Implementation", "An embedded simulation terminal mimicking a server terminal node (`system@growthlytics:~#`)."),
    ("Dynamic Mechanics", "Users type issues (e.g. 'manual entry mistakes') and the model generates operational solutions on-screen."),
    ("Strategic Impact", "Highlights the bridging of LLM APIs with real-world small business workflows.")
]

for title, body in points:
    p = tf.add_paragraph()
    p.text = f"• {title}: "
    p.font.name = 'Arial'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    run = p.add_run()
    run.text = body
    run.font.name = 'Arial'
    run.font.size = Pt(13)
    run.font.bold = False
    run.font.color.rgb = GRAY
    p.space_after = Pt(14)

ai_img = cropped_paths.get("04_ai_advisor")
if ai_img:
    slide.shapes.add_picture(ai_img, Inches(7.2), Inches(1.6), width=Inches(5.6))

# Slide 5: Featured Project: SkillsElevate
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_slide_header(slide, "Pinned Project: SkillsElevate Portal", "04 -- FEATURED PRODUCT")

tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.8), Inches(5.2))
tf = tb.text_frame
tf.word_wrap = True

pts = [
    ("Core Description", "An educational upskilling portal engineered to help working professionals and students upgrade their capabilities."),
    ("AI Integration", "Embedded Google Gemini API curates customized learning trajectories and maps real-time industry trends."),
    ("Technical Stack", "React, TypeScript, Google Gemini SDK, Node.js, and Vite build wrapper."),
    ("Key Features", "Generates interactive roadmaps, defines daily learning checkpoints, and supports multi-language learning tracks.")
]

for title, body in pts:
    p = tf.add_paragraph()
    p.text = f"• {title}: "
    p.font.name = 'Arial'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = VIOLET
    
    run = p.add_run()
    run.text = body
    run.font.name = 'Arial'
    run.font.size = Pt(13)
    run.font.bold = False
    run.font.color.rgb = GRAY
    p.space_after = Pt(16)

proj_img = cropped_paths.get("05_projects")
if proj_img:
    slide.shapes.add_picture(proj_img, Inches(7.5), Inches(1.6), width=Inches(5.3))

# Slide 6: Live Web App: Growthlytics & Momotantra
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_slide_header(slide, "Growthlytics & Momotantra Deployed Apps", "05 -- CLIENT CASE STUDIES")

tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.2), Inches(5.2))
tf = tb.text_frame
tf.word_wrap = True

p_g = tf.paragraphs[0]
p_g.text = "1. GROWTHLYTICS (Flagship Agency)"
p_g.font.name = 'Arial'
p_g.font.size = Pt(16)
p_g.font.bold = True
p_g.font.color.rgb = GOLD
p_g.space_after = Pt(4)

g_desc = tf.add_paragraph()
g_desc.text = "• Flagship agency portal offering local business scaling, AI diagnostics, and custom SEO management dashboards.\n• Purpose: Bridge the digital divide for small vendors in West Bengal.\n• Live Link: growthlytics-s6rd.vercel.app"
g_desc.font.name = 'Arial'
g_desc.font.size = Pt(12)
g_desc.font.color.rgb = GRAY
g_desc.space_after = Pt(18)

p_m = tf.add_paragraph()
p_m.text = "2. MOMOTANTRA (Digital Menu & Billing)"
p_m.font.name = 'Arial'
p_m.font.size = Pt(16)
p_m.font.bold = True
p_m.font.color.rgb = TEAL
p_m.space_after = Pt(4)

m_desc = tf.add_paragraph()
m_desc.text = "• Premium interactive digital menu card and checkout panel designed for street-side food stalls and momo vendors.\n• Purpose: Accelerate checkout speeds, automate billing invoices, and minimize customer queues.\n• Live Link: momotantra-app-a3mw.vercel.app"
m_desc.font.name = 'Arial'
m_desc.font.size = Pt(12)
m_desc.font.color.rgb = GRAY

# Add Growthlytics & Momotantra images side-by-side
g_screenshot = sub_apps.get("growthlytics")
if g_screenshot and os.path.exists(g_screenshot):
    slide.shapes.add_picture(g_screenshot, Inches(7.0), Inches(1.6), width=Inches(2.8))

m_screenshot = sub_apps.get("momotantra")
if m_screenshot and os.path.exists(m_screenshot):
    slide.shapes.add_picture(m_screenshot, Inches(10.1), Inches(1.6), width=Inches(2.8))

# Slide 7: Live Web App: Rever Cafe & RetailFlow AI
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_slide_header(slide, "Rever Cafe & RetailFlow Deployed Apps", "06 -- DIGITAL UTILITIES")

tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.2), Inches(5.2))
tf = tb.text_frame
tf.word_wrap = True

p_r = tf.paragraphs[0]
p_r.text = "3. REVER GLASS HOUSE CAFE"
p_r.font.name = 'Arial'
p_r.font.size = Pt(16)
p_r.font.bold = True
p_r.font.color.rgb = VIOLET
p_r.space_after = Pt(4)

r_desc = tf.add_paragraph()
r_desc.text = "• Luxury-themed sensory landing page and booking interface built for a modern aesthetic cafe network.\n• Features: Interactive booking scheduler, animated menus, and integrated feedback portals.\n• Live Link: revercafe.vercel.app"
r_desc.font.name = 'Arial'
r_desc.font.size = Pt(12)
r_desc.font.color.rgb = GRAY
r_desc.space_after = Pt(18)

p_f = tf.add_paragraph()
p_f.text = "4. RETAILFLOW AI & TECHAI"
p_f.font.name = 'Arial'
p_f.font.size = Pt(16)
p_f.font.bold = True
p_f.font.color.rgb = TEAL
p_f.space_after = Pt(4)

f_desc = tf.add_paragraph()
f_desc.text = "• Retail inventory modeling and sales tracking application utilizing time-series forecasting modules to predict stock demand and run custom ledger reports.\n• TechAI: Secure OTP authentication dashboard with AI coding agents.\n• Live Link: retailflow-ai-opal.vercel.app"
f_desc.font.name = 'Arial'
f_desc.font.size = Pt(12)
f_desc.font.color.rgb = GRAY

# Add Rever Cafe & RetailFlow screenshots side-by-side
r_screenshot = sub_apps.get("revercafe")
if r_screenshot and os.path.exists(r_screenshot):
    slide.shapes.add_picture(r_screenshot, Inches(7.0), Inches(1.6), width=Inches(2.8))

f_screenshot = sub_apps.get("retailflow")
if f_screenshot and os.path.exists(f_screenshot):
    slide.shapes.add_picture(f_screenshot, Inches(10.1), Inches(1.6), width=Inches(2.8))

# Slide 8: Agency Hub & Vision
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_slide_header(slide, "Growthlytics Vision & Expansion Hub", "07 -- STRATEGIC ROADMAP")

tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.5), Inches(5.2))
tf = tb.text_frame
tf.word_wrap = True

vpoints = [
    ("Mission Statement", "Connecting offline to online. Democratize enterprise-grade digital tools (AI, SEO, databases) for local micro-vendors and MSMEs in West Bengal."),
    ("Hyper-Local SEO Node", "Assisting vendors to expand local organic search discovery by optimized Google Business listings and custom SEO landing sheets."),
    ("AI-Assisted Sales Modules", "Deploying cost-effective database structures (Advanced Excel, SQL tables) so street-side stores can forecast weekly inventory requirements."),
    ("Hub Contacts", "Active Headquarters at 11/1 Arabinda Road, Konnagar, West Bengal, India. 24/7 hotline support (+91 98361 48511).")
]

for title, body in vpoints:
    p = tf.add_paragraph()
    p.text = f"• {title}: "
    p.font.name = 'Arial'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = GOLD
    
    run = p.add_run()
    run.text = body
    run.font.name = 'Arial'
    run.font.size = Pt(13)
    run.font.bold = False
    run.font.color.rgb = GRAY
    p.space_after = Pt(14)

hub_img = cropped_paths.get("06_agency_hub")
if hub_img:
    slide.shapes.add_picture(hub_img, Inches(7.3), Inches(1.6), width=Inches(5.5))

# Slide 9: Resume & Career Nodes
slide = prs.slides.add_slide(slide_layout)
apply_background(slide)
add_slide_header(slide, "Interactive Resume & Professional Credentials", "08 -- CAREER TRACK")

tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.8), Inches(5.2))
tf = tb.text_frame
tf.word_wrap = True

rpts = [
    ("Work Experience", "Founder of Growthlytics (Present); Customer Service Associate & Data Coordinator at Tech Mahindra (Jun 2025 - Feb 2026 - handled daily SLAs and weekly reports); Business Analytics Intern at Bharat Intern."),
    ("Key Certifications", "Google Digital Marketing & E-Commerce, Yale University Capital Markets, PwC Power BI Case, J.P. Morgan Global Finance, and J.P. Morgan Corporate Analyst Development."),
    ("Academic Credentials", "Bachelor of Arts (2021 - 2024) from MGKVP University with CGPA of 6.62 / 10; commerce stream background (Class XII, 69%)."),
    ("Operations Record", "Coordinated outreach and data logging for 100+ entries as Campus Ambassador for E-Cell IIT Bombay.")
]

for title, body in rpts:
    p = tf.add_paragraph()
    p.text = f"• {title}: "
    p.font.name = 'Arial'
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = TEAL
    
    run = p.add_run()
    run.text = body
    run.font.name = 'Arial'
    run.font.size = Pt(13)
    run.font.bold = False
    run.font.color.rgb = GRAY
    p.space_after = Pt(14)

res_img = cropped_paths.get("07_resume")
if res_img:
    slide.shapes.add_picture(res_img, Inches(7.5), Inches(1.5), width=Inches(5.3))

# Save the presentation
prs.save(OUTPUT_PPTX_PATH)
print(f"Presentation saved successfully to: {OUTPUT_PPTX_PATH}")


# ----------------- GENERATING ANIMATED GIF (WALKTHROUGH VIDEO) -----------------
print("Generating Portfolio Walkthrough Animated GIF (Video)...")
# We will create an animated GIF scrolling through the full screenshot image.
# We will slide a viewport of 1280x720 from top to bottom in increments of 60px.
# This creates a perfect simulated "screencast video" of the website!
gif_slides = []

# Crop settings for GIF:
viewport_w = 1280
viewport_h = 720
scroll_step = 65  # scroll step in pixels
max_y = height - viewport_h

# Let's crop frames
current_y = 0
frame_count = 0
while current_y <= max_y and frame_count < 100:  # Cap at 100 frames to avoid memory exhaustion
    # Crop the viewport
    frame = img.crop((0, current_y, viewport_w, current_y + viewport_h))
    
    # Resize frame slightly to make GIF smaller in file size and load faster
    frame_resized = frame.resize((854, 480), Image.Resampling.LANCZOS)
    gif_slides.append(frame_resized)
    
    # Scroll slower on key sections
    is_near_section = False
    for section_name, bbox in crops.items():
        sect_top = bbox[1]
        if abs(current_y - sect_top) < 100:
            is_near_section = True
            break
            
    if is_near_section:
        # Repeat the frame to create a pause/delay on key sections
        for _ in range(3):
            gif_slides.append(frame_resized)
            
    current_y += scroll_step
    frame_count += 1

# Append sub-app screenshots to the end of the GIF
for app_name, app_path in sub_apps.items():
    if os.path.exists(app_path):
        app_img = Image.open(app_path)
        # Resize to fit the GIF resolution (854x480)
        app_frame = app_img.resize((854, 480), Image.Resampling.LANCZOS)
        
        # Add a title text overlay to the app screenshot
        draw = ImageDraw.Draw(app_frame)
        # Use default font or simple styling
        draw.rectangle([(10, 10), (280, 50)], fill=(10, 10, 15, 200))
        draw.text((20, 20), f"APP: {app_name.upper()}", fill=(201, 168, 76))
        
        # Pause on each app screenshot
        for _ in range(12):
            gif_slides.append(app_frame)

if gif_slides:
    # Save animated GIF
    # duration is in milliseconds per frame (e.g. 100ms = 10fps)
    gif_slides[0].save(
        OUTPUT_GIF_PATH,
        save_all=True,
        append_images=gif_slides[1:],
        optimize=True,
        duration=120,
        loop=0
    )
    print(f"Animated GIF walkthrough saved successfully to: {OUTPUT_GIF_PATH}")
else:
    print("[ERROR] No frames generated for animated GIF.")

print("All tasks completed successfully!")
